"""
Сборка презентации (.pptx) из графиков и цифр в БД.

Презентация рассчитана на чтение без докладчика: всё, что обычно
проговаривают голосом, вынесено на слайд текстом. Отсюда плотность —
это slidedoc, а не проекционная презентация.

Числа берутся из Postgres при каждой сборке (см. charts/queries.py).
Формулировки выводов — константы в этом файле: они не считаются,
а пишутся человеком.

Семь слайдов, план — docs/presentation_plan.md. Детали, не поместившиеся
на слайды, — docs/findings.md.

Запуск (из корня проекта, после etl/main.py и charts/main.py):
    python deck.py

Зависимости: python-pptx (см. requirements.txt).
"""

import sys
from pathlib import Path

# charts/ — не пакет, кладём в путь импорта, чтобы переиспользовать
# уже написанные запросы (тот же приём, что в verify.py).
sys.path.insert(0, str(Path(__file__).resolve().parent / "charts"))

import psycopg2  # noqa: E402
import queries  # noqa: E402
from config import DB_CONFIG  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

ROOT = Path(__file__).resolve().parent
CHARTS = ROOT / "charts" / "output"
OUT_FILE = ROOT / "portfolio_audit.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)
FULL_W = Inches(12.23)

INK = RGBColor(0x0B, 0x0B, 0x0B)
INK_SOFT = RGBColor(0x52, 0x51, 0x4E)
INK_MUTE = RGBColor(0x89, 0x87, 0x81)
ACCENT = RGBColor(0x2A, 0x78, 0xD6)
ACCENT_BG = RGBColor(0xEC, 0xF3, 0xFD)
WARN = RGBColor(0xD9, 0x53, 0x1E)
WARN_BG = RGBColor(0xFD, 0xF0, 0xEA)
PAPER = RGBColor(0xF6, 0xF5, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Arial"


# ----------------------------------------------------------- примитивы

def _text(slide, left, top, width, height, blocks, align=PP_ALIGN.LEFT,
          spacing=Pt(4)):
    """blocks — список кортежей (текст, кегль, цвет, жирный)."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (txt, size, color, bold) in enumerate(blocks):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.space_after = spacing
        run = para.add_run()
        run.text = txt
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT
    return box


def _one(slide, left, top, width, height, txt, size, color=INK, bold=False):
    return _text(slide, left, top, width, height, [(txt, size, color, bold)])


def _eyebrow_title(slide, eyebrow, title):
    _one(slide, MARGIN, Inches(0.28), FULL_W, Inches(0.24), eyebrow, 10, INK_MUTE)
    _one(slide, MARGIN, Inches(0.56), FULL_W, Inches(0.72), title, 24, INK, bold=True)


def _panel(slide, left, top, width, height, fill, edge=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.05
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if edge is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = edge
        box.line.width = Pt(1)
    box.shadow.inherit = False
    return box


def _picture(slide, name, left, top, height):
    path = CHARTS / name
    if not path.exists():
        sys.exit(f"Не найден график: {path}\nЗапускали charts/main.py?")
    return slide.shapes.add_picture(str(path), left, top, height=height)


def _table(slide, left, top, width, header, rows_data, col_w,
           row_h=Inches(0.3), size=10, highlight=None):
    n_rows, n_cols = len(rows_data) + 1, len(header)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_h * n_rows)
    tbl = shape.table
    tbl.first_row = True
    for i, frac in enumerate(col_w):
        tbl.columns[i].width = Emu(int(width * frac))
    for i in range(n_rows):
        tbl.rows[i].height = row_h

    def fill(cell, txt, bold, color, right):
        cell.text = str(txt)
        cell.margin_left = cell.margin_right = Inches(0.06)
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT if right else PP_ALIGN.LEFT
        r = p.runs[0]
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FONT

    for j, txt in enumerate(header):
        fill(tbl.cell(0, j), txt, True, WHITE, j > 1)
    for i, row in enumerate(rows_data, start=1):
        accent = highlight is not None and i - 1 == highlight
        for j, txt in enumerate(row):
            fill(tbl.cell(i, j), txt, accent, ACCENT if accent else INK, j > 1)
    return shape


def _source(slide, txt):
    _one(slide, MARGIN, Inches(7.14), FULL_W, Inches(0.2), txt, 8, INK_MUTE)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# -------------------------------------------------------------- слайды

def slide_1(prs, kpi, cost):
    s = _blank(prs)
    _one(s, MARGIN, Inches(0.26), Inches(9.0), Inches(0.24),
         "АУДИТ ПОРТФЕЛЯ ИТ-ПРОЕКТОВ   ·   100 ПРОЕКТОВ   ·   ДАННЫЕ ЗА 2025 ГОД",
         10, INK_MUTE)
    _one(s, MARGIN, Inches(0.54), FULL_W, Inches(0.74),
         "Портфель выполняется, но треть его стоимости — в проектах, сорвавших срок",
         24, INK, bold=True)

    _panel(s, MARGIN, Inches(1.38), FULL_W, Inches(1.0), PAPER)
    _text(s, Inches(0.75), Inches(1.5), Inches(11.83), Inches(0.85), [
        ("Три выгрузки из независимых систем: PMO (сроки и задачи), финансовый учёт "
         "(платежи), ИТ-департамент (закупки). 600 задач, 768 платежей, 355 позиций "
         "оборудования. Связаны по коду проекта после нормализации — в источнике коды "
         "записаны в двух форматах.", 9, INK_SOFT, False),
        ("Допущения: курс фиксированный, 1 USD = 90 ₽, 1 EUR = 100 ₽. Категория проекта "
         "определяется составом закупок. Четыре проекта с незакрытой последней задачей "
         "исключены из расчёта надёжности. Полный разбор качества данных — docs/findings.md.",
         9, INK_MUTE, False),
    ], spacing=Pt(3))

    for i, (value, label) in enumerate(kpi):
        left = MARGIN + Inches(i * 3.06)
        _one(s, left, Inches(2.6), Inches(2.9), Inches(0.48), value, 26, ACCENT, bold=True)
        _one(s, left, Inches(3.08), Inches(2.9), Inches(0.26), label, 10, INK_SOFT)
    _one(s, MARGIN, Inches(3.44), FULL_W, Inches(0.24),
         "По этим показателям портфель выглядит здоровым.", 11, INK_MUTE)

    _panel(s, MARGIN, Inches(3.82), FULL_W, Inches(1.04), WARN_BG, WARN)
    _text(s, Inches(0.75), Inches(3.95), Inches(11.83), Inches(0.88), [
        (f"{cost['late_projects']} проекта не уложились в срок. На них приходится "
         f"{cost['late_bln']} млрд ₽ — {cost['late_pct']}% стоимости портфеля.",
         15, WARN, True),
        ("Средний срыв по портфелю — 2,4 дня. В днях проблема выглядит незначительной. "
         "В деньгах — нет.", 10, INK_SOFT, False),
    ], spacing=Pt(4))

    _one(s, MARGIN, Inches(5.06), FULL_W, Inches(0.26),
         "Что из этого следует", 12, INK, bold=True)
    lines = [
        ("1.   Сроки занижают при оценке, а не срывают при исполнении.",
         "За план выходят все шесть этапов, коэффициент факт/план — 1,076.        → слайды 3–4"),
        ("2.   Поправка нужна разная для разных типов проектов, сейчас она одна на всех.",
         "Занижение 6,8% против 10,3% при одинаковом буфере.        → слайд 5"),
        ("3.   96% стоимости портфеля зависит от курса, который никто не пересчитывал.",
         "Сдвиг курса на 10% — это 1,15 млрд ₽.        → слайд 6"),
        ("4.   Часть отчётности нельзя проверить.",
         "Источники подставляют заглушки вместо пустых значений.        → слайд 7"),
    ]
    top = Inches(5.38)
    for head, body in lines:
        _one(s, MARGIN, top, FULL_W, Inches(0.22), head, 11, INK, bold=True)
        _one(s, Inches(0.88), top + Inches(0.2), Inches(11.9), Inches(0.22),
             body, 9, INK_MUTE)
        top += Inches(0.42)

    _source(s, "Источник: sql/analysis.sql (2.1), sql/bonus_analysis.sql (7).")


def slide_2(prs, top5, relative, cost):
    s = _blank(prs)
    _eyebrow_title(s, "СТОИМОСТЬ ПРОСРОЧКИ",
                   "Треть стоимости портфеля лежит в проектах, не уложившихся в срок")

    _one(s, MARGIN, Inches(1.36), Inches(5.5), Inches(0.24),
         "68 проектов из 100 укладываются в 0–3 дня — проблема в хвосте", 11, INK, bold=True)
    _picture(s, "overrun_histogram.png", MARGIN, Inches(1.66), Inches(2.2))
    _one(s, MARGIN, Inches(3.94), Inches(5.5), Inches(0.44),
         "Столбец на нуле — 58 проектов, сданных день в день. Хвост — 23 проекта "
         "со срывом 6–19 дней.", 9, INK_MUTE)

    left = Inches(6.3)
    _one(s, left, Inches(1.36), Inches(6.48), Inches(0.24),
         "Пять проектов с наибольшим срывом", 11, INK, bold=True)
    header = ["Проект", "Руководитель", "Срыв", "% плана", "Бюджет"]
    rows_data = [
        (f"{code}  {name}", pm, f"{days} дн.", f"{pct}%", f"{budget} млн")
        for code, name, pm, days, pct, budget, _open in top5
    ]
    _table(s, left, Inches(1.66), Inches(6.48), header, rows_data,
           col_w=(0.40, 0.22, 0.12, 0.13, 0.13), size=9, highlight=0)

    worst_abs, worst_rel = top5[0], relative[0]
    _text(s, left, Inches(3.76), Inches(6.48), Inches(1.5), [
        (f"Больше всех дней потерял {worst_abs[0]} «{worst_abs[1]}» — {worst_abs[3]} дней. "
         f"Но при плановой длительности 206 дней это всего {worst_abs[4]}%.", 10, INK, True),
        (f"Хуже всех держал собственный график {worst_rel[0]} «{worst_rel[1]}»: те же "
         f"{worst_rel[2]} дней, но при плане в {worst_rel[3]} — это {worst_rel[4]}%.",
         10, INK_SOFT, False),
        ("Рейтинги по абсолютному и относительному срыву совпадают на три проекта из пяти. "
         "«Потерял больше всех дней» и «хуже всех держал план» — разные проекты.",
         10, INK_SOFT, False),
    ], spacing=Pt(5))

    _panel(s, MARGIN, Inches(5.44), FULL_W, Inches(1.0), ACCENT_BG)
    _text(s, Inches(0.75), Inches(5.57), Inches(11.83), Inches(0.84), [
        (f"{cost['late_projects']} проекта   ·   {cost['late_bln']} млрд ₽   ·   "
         f"{cost['late_pct']}% портфеля", 17, ACCENT, True),
        (f"Хвост из {cost['tail_projects']} проектов со срывом от 6 дней держит "
         f"{cost['tail_bln']} млрд ₽ — {cost['tail_pct']}% портфеля.", 10, INK_SOFT, False),
    ], spacing=Pt(3))

    _one(s, MARGIN, Inches(6.6), FULL_W, Inches(0.22),
         "У четырёх проектов последняя задача не закрыта, по ним срыв занижен. "
         "Разбор — docs/findings.md.", 9, INK_MUTE)
    _source(s, "Источник: sql/analysis.sql (1.2), sql/bonus_analysis.sql (5, 6, 7).")


def slide_3(prs, stages):
    s = _blank(prs)
    _eyebrow_title(s, "ГДЕ ТЕРЯЕТСЯ ВРЕМЯ",
                   "Опаздывает не какой-то один этап — за план выходят все шесть")

    _one(s, MARGIN, Inches(1.36), Inches(6.5), Inches(0.24),
         "Ни один этап в среднем не укладывается в план", 11, INK, bold=True)
    _picture(s, "stage_benchmark_bar.png", MARGIN, Inches(1.66), Inches(2.3))

    left = Inches(7.3)
    _one(s, left, Inches(1.36), Inches(5.48), Inches(0.24),
         "Плановая и фактическая длительность, дней", 11, INK, bold=True)
    header = ["Этап", "План", "Факт", "Отклонение"]
    rows_data = [(name, f"{plan}", f"{fact}", f"+{pct}%")
                 for name, _n, plan, fact, _r, pct in stages]
    worst = max(range(len(stages)), key=lambda i: stages[i][5])
    _table(s, left, Inches(1.66), Inches(5.48), header, rows_data,
           col_w=(0.46, 0.16, 0.16, 0.22), size=9, highlight=worst)
    _one(s, left, Inches(3.86), Inches(5.48), Inches(0.22),
         "По 585 задачам из 600, где заполнены обе фактические даты.", 9, INK_MUTE)

    _panel(s, MARGIN, Inches(4.24), Inches(6.5), Inches(0.6), PAPER)
    _one(s, Inches(0.72), Inches(4.36), Inches(6.16), Inches(0.4),
         "Разница между худшим и лучшим этапом — 1,2 дня. Проблема не локализована "
         "в конкретной работе.", 9, INK_SOFT)

    _text(s, MARGIN, Inches(5.06), FULL_W, Inches(1.1), [
        ("Разработка модуля и тестирование срываются сильнее прочих, но ненамного: "
         "+3,8 и +3,7 дня против +2,6 у анализа требований. Все шесть значений положительные.",
         11, INK_SOFT, False),
        ("Если бы дело было в конкретных работах, отклонения распределялись бы неравномерно: "
         "часть этапов укладывалась бы в план, часть выбивалась. Здесь выбиваются все. "
         "Значит, смещена сама методика оценки, а не исполнение отдельных задач.",
         11, INK, True),
    ], spacing=Pt(5))

    _one(s, MARGIN, Inches(6.46), FULL_W, Inches(0.5),
         "Проверенная гипотеза: срыв накапливается каскадом, ранние этапы задерживают поздние. "
         "Средний срыв по этапам в порядке выполнения — 2,56 / 3,31 / 3,80 / 3,71 / 3,42 / 2,86. "
         "Кривая ровная, к концу проекта не растёт. Не подтвердилась.", 9, INK_MUTE)
    _source(s, "Источник: sql/analysis.sql (1.1), sql/bonus_analysis.sql (1, 9).")


def slide_4(prs, balance, stages, leak):
    s = _blank(prs)
    _eyebrow_title(s, "МЕХАНИКА СРЫВА",
                   "Оценка занижена на 8%, а буфер в плане рассчитан впритык")

    _one(s, MARGIN, Inches(1.34), Inches(6.0), Inches(0.24),
         f"Из {balance['accumulated']} дня накопленных задержек до срока проекта "
         f"доходит {balance['overrun']}", 11, INK, bold=True)
    _picture(s, "buffer_waterfall.png", MARGIN, Inches(1.62), Inches(2.6))

    left = Inches(6.85)
    lo = min(stages, key=lambda r: r[5])
    hi = max(stages, key=lambda r: r[5])
    _text(s, left, Inches(1.34), Inches(5.93), Inches(1.6), [
        ("Откуда 8%", 12, INK, True),
        (f"Факт длиннее плана на каждом этапе: от +{lo[5]}% ({lo[0].lower()}) "
         f"до +{hi[5]}% ({hi[0].lower()}).", 10, INK_SOFT, False),
        ("По 585 измеренным задачам средний коэффициент факт/план — 1,076. Оценки занижены "
         "примерно на 8%, равномерно по всем видам работ.", 10, INK_SOFT, False),
        ("Считается как (fact_end − fact_start) / (plan_end − plan_start) по каждой задаче, "
         "где заполнены обе фактические даты.", 9, INK_MUTE, False),
    ], spacing=Pt(4))

    _text(s, left, Inches(3.18), Inches(5.93), Inches(1.55), [
        ("Откуда буфер", 12, INK, True),
        (f"Следующий этап стартует не сразу после планового конца предыдущего. Медиана "
         f"зазора — ровно {balance['gap_per_transition']} дня, из {leak['n_transitions']} "
         f"переходов ни одного перекрытия. В сумме это {balance['buffer']} дней запаса.",
         10, INK_SOFT, False),
        (f"Но буфер расходуется по каждому переходу отдельно, не общей суммой: на "
         f"{leak['leak_pct']}% переходов ({leak['leak_transitions']} из "
         f"{leak['n_transitions']}) срыв конкретного этапа больше зазора именно перед ним, "
         f"и тогда в среднем утекает ещё {leak['avg_leak_days']} дня дальше по проекту.",
         10, INK_SOFT, False),
    ], spacing=Pt(4))

    _panel(s, MARGIN, Inches(4.94), FULL_W, Inches(0.6), ACCENT_BG, ACCENT)
    _one(s, Inches(0.75), Inches(5.06), Inches(11.83), Inches(0.4),
         "Буфер есть, и он работает: гасит 85% отклонений в среднем по проекту. "
         "Проблема не в его отсутствии, а в том, что он рассчитан впритык и неравномерно.",
         13, ACCENT, bold=True)

    _text(s, MARGIN, Inches(5.7), FULL_W, Inches(1.3), [
        ("Это описанный эффект: оценивая собственную работу, люди смотрят «изнутри» — "
         "на конкретный план — и недооценивают, чем кончались похожие работы раньше. "
         "Систематическое занижение на единицы процентов при этом типично. Стандартное "
         "решение — считать поправку по прошлым проектам, а не по экспертному суждению.",
         10, INK_SOFT, False),
        ("Методу нужна база сопоставимых завершённых проектов, и обычно её нет. Здесь она "
         "есть: 100 проектов, 6 стандартных этапов, 585 измеренных задач. Коэффициент "
         "не надо угадывать — он посчитан.", 10, INK, True),
    ], spacing=Pt(5))

    _source(s, "Источник: sql/analysis.sql (1.1), sql/bonus_analysis.sql (1, 9, 10, 12).")


def slide_5(prs, cats):
    s = _blank(prs)
    _eyebrow_title(s, "КАТЕГОРИИ ПРОЕКТОВ",
                   "Буфер один на всех, а ошибка оценки у разных проектов разная")

    header = ["Категория", "Занижение", "Буфер", "Средний срыв", "Опозданий"]
    rows_data = [(name, f"{pct}%", "15 дн.", f"+{overrun} дн.", f"{late}%")
                 for name, _n, pct, overrun, late in cats]
    worst = max(range(len(cats)), key=lambda i: cats[i][2])
    _table(s, MARGIN, Inches(1.46), Inches(7.05), header, rows_data,
           col_w=(0.34, 0.17, 0.14, 0.20, 0.15), row_h=Inches(0.33), size=11,
           highlight=worst)
    _one(s, MARGIN, Inches(2.56), Inches(7.05), Inches(0.22),
         "Колонка «буфер» одинаковая — в этом и суть.", 9, INK_MUTE)

    _picture(s, "category_comparison.png", Inches(7.9), Inches(1.46), Inches(1.85))
    _one(s, Inches(7.9), Inches(3.4), Inches(4.88), Inches(0.44),
         "Офисные проекты дешевле на треть, но опаздывают вдвое сильнее.", 9, INK_MUTE)

    _text(s, MARGIN, Inches(3.06), Inches(7.05), Inches(1.9), [
        ("Офисно-коммерческие проекты опаздывают не чаще технологических: 34,8% против "
         "32,5%, разница в пределах погрешности. Они опаздывают сильнее — вдвое.",
         11, INK_SOFT, False),
        ("Причина в первой колонке: ошибка оценки у них в полтора раза больше, "
         "а поправка применяется одинаковая.", 11, INK_SOFT, False),
        ("Единый буфер — неправильный инструмент. Калибровать нужно по категории закупки, "
         "а не одной цифрой на весь портфель.",
         12, INK, True),
    ], spacing=Pt(5))

    _panel(s, MARGIN, Inches(5.3), FULL_W, Inches(1.55), PAPER)
    _text(s, Inches(0.75), Inches(5.44), Inches(11.83), Inches(1.35), [
        ("Как определяется категория", 11, INK, True),
        ("Тип проекта — банк, склад, ЦОД, ритейл и т.д. — с категорией не связан: один и тот "
         "же тип встречается в обеих группах. Категория отражает состав закупки, а не отрасль "
         "проекта.", 10, INK_SOFT, False),
        ("Технологические — проекты, где закупались серверы или системные блоки. Монитор "
         "в правило не входит: встречается в обеих группах примерно одинаково — 42 из 77 "
         "и 15 из 23.", 10, INK_SOFT, False),
        ("15 проектов лежат на границе — от классификации зависит коэффициент поправки. "
         "Рекомендация — закрепить правило в регламенте.", 10, INK, True),
    ], spacing=Pt(3))

    _source(s, "Источник: sql/analysis.sql (3.1), sql/bonus_analysis.sql (1, 11).")


def slide_6(prs, concentration, pm_top3):
    s = _blank(prs)
    _eyebrow_title(s, "ФИНАНСОВЫЙ РИСК",
                   "Второй крупный риск — курс валют, и его никто не считал")

    _one(s, MARGIN, Inches(1.34), Inches(4.2), Inches(0.24),
         "96% стоимости — платежи в валюте", 11, INK, bold=True)
    _picture(s, "currency_donut.png", MARGIN, Inches(1.62), Inches(2.3))
    _one(s, MARGIN, Inches(4.0), Inches(4.2), Inches(0.6),
         "В рублях 80% платежей, но лишь 3,8% стоимости. Рублёвые платежи многочисленные "
         "и мелкие, валютные — редкие и крупные.", 9, INK_MUTE)

    left = Inches(5.1)
    _panel(s, left, Inches(1.34), Inches(7.68), Inches(2.2), ACCENT_BG, ACCENT)
    _text(s, left + Inches(0.22), Inches(1.46), Inches(7.24), Inches(2.0), [
        ("Пересчёт сделан по фиксированному курсу: 1 USD = 90 ₽, 1 EUR = 100 ₽.",
         10, INK_SOFT, False),
        ("Сдвиг курса на 5% меняет оценку портфеля на 0,58 млрд ₽", 10, INK_SOFT, False),
        ("Сдвиг курса на 10% — на 1,15 млрд ₽", 14, ACCENT, True),
        ("Сдвиг курса на 20% — на 2,31 млрд ₽", 10, INK_SOFT, False),
        ("12,0 млрд — не измеренная величина, а расчёт при заданном курсе. Показывать её "
         "одним числом некорректно.", 11, INK, True),
    ], spacing=Pt(4))

    _one(s, left, Inches(3.72), Inches(7.68), Inches(0.24),
         "Концентрация бюджета", 12, INK, bold=True)
    pm_line = "     ".join(f"{name} — {bln} млрд ₽" for name, bln in pm_top3)
    _one(s, left, Inches(4.0), Inches(7.68), Inches(0.24), pm_line, 10, INK_SOFT)
    conc = "        ".join(f"{seg}: {bln} млрд ₽ ({pct}% портфеля)"
                           for seg, bln, pct in concentration)
    _one(s, left, Inches(4.3), Inches(7.68), Inches(0.3), conc, 11, ACCENT, bold=True)

    _panel(s, MARGIN, Inches(4.94), FULL_W, Inches(1.6), PAPER)
    _text(s, Inches(0.75), Inches(5.08), Inches(11.83), Inches(1.4), [
        ("Проверенная гипотеза: руководители с большей нагрузкой чаще срывают сроки",
         11, INK, True),
        ("Связь числа проектов с долей опозданий — −0,07, то есть её нет. Васильев ведёт "
         "15 проектов и опаздывает в 13,3% случаев; Смирнов ведёт 8 и опаздывает в 37,5%. "
         "Гипотеза не подтвердилась.", 10, INK_SOFT, False),
        ("Кроме того, на восьми руководителях по 8–16 проектов у каждого разброс 13–56% "
         "может объясняться случайностью. Кадровых выводов на такой выборке не делаем.",
         10, INK_SOFT, False),
    ], spacing=Pt(3))

    _one(s, MARGIN, Inches(6.66), FULL_W, Inches(0.22),
         "Двадцать крупнейших проектов из ста дают 47,5% стоимости портфеля.", 9, INK_MUTE)
    _source(s, "Источник: sql/analysis.sql (2.1, 2.2), sql/bonus_analysis.sql (2, 3, 8).")


def slide_7(prs):
    s = _blank(prs)
    _eyebrow_title(s, "РЕКОМЕНДАЦИИ", "Что делать и по какой метрике проверять")

    cards = [
        ("1.   Считать поправку к оценке по прошлым проектам, а не экспертно",
         "Коэффициент измерен по 585 задачам: 1,08 в среднем, 1,24 на уровне 80% задач, "
         "1,48 на уровне 90%. Выбор уровня — вопрос того, какую долю проектов допустимо "
         "сдавать с опозданием. Меняется коэффициент в методике, а не работа людей.",
         "Метрики нет. Ввести: коэффициент факт/план по этапу, помесячно. "
         "Цель — 1,0. Сегодня 1,076."),
        ("2.   Развести поправку по типам проектов",
         "Технологическим — 6,8%, офисно-коммерческим — 10,3%. Сейчас обе группы получают "
         "одинаковый буфер в 15 дней.",
         "Метрика: тот же коэффициент в разрезе категорий."),
        ("3.   Расширить буфер с 15 до 17 дней либо следить за его расходом",
         "15 дней покрывают половину проектов, 17 — девяносто пять процентов.",
         "Метрики нет. Ввести: доля израсходованного буфера по проекту — позволяет "
         "вмешаться до срыва, а не узнать о нём после."),
        ("4.   Перестать показывать стоимость портфеля одним числом",
         "96% зависит от курса. Либо фиксировать курс в контрактах, либо публиковать вилку.",
         "Метрика: доля валютных обязательств и чувствительность к сдвигу на 10%."),
        ("5.   Убрать заглушки из форм ввода",
         "Дефекты систематические: битые коды 4,8% строк, суммы текстом 3,9%, дата-заглушка "
         "в платежах 2,6%, задачи без даты начала 2,5%, даты-заглушки в сроках 1,7%, плюс "
         "6 задач с нарушенным порядком выполнения. Одно и то же значение повторяется "
         "во всех затронутых строках — это поведение форм ввода, а не ошибки людей.",
         "Метрика: доля строк с дефектами по источнику, помесячно. "
         "Разбор — docs/findings.md."),
    ]
    top = Inches(1.38)
    for head, body, metric in cards:
        _one(s, MARGIN, top, Inches(8.5), Inches(0.22), head, 11, INK, bold=True)
        _one(s, Inches(0.88), top + Inches(0.21), Inches(8.1), Inches(0.5), body, 9, INK_SOFT)
        _one(s, Inches(9.28), top + Inches(0.02), Inches(3.5), Inches(0.66), metric, 9, ACCENT)
        top += Inches(1.0)

    _panel(s, MARGIN, Inches(6.46), FULL_W, Inches(0.66), PAPER)
    _text(s, Inches(0.75), Inches(6.56), Inches(11.83), Inches(0.56), [
        ("Чего не хватает, чтобы проверить остальное", 10, INK, True),
        ("Перегрузка команды — нет трудозатрат и состава команды.      "
         "Внешние блокировки — нет причины задержки.      "
         "Сдвиг дедлайнов задним числом — нет истории версий плана.", 9, INK_SOFT, False),
    ], spacing=Pt(2))


# --------------------------------------------------------------- запуск

KPI = [
    ("100", "проектов в портфеле"),
    ("12,0 млрд ₽", "общая стоимость"),
    ("120,0 млн ₽", "средний проект"),
    ("98,3%", "задач закрыто"),
]


def main():
    if not CHARTS.exists():
        sys.exit(f"Нет каталога с графиками: {CHARTS}\nЗапустите сначала charts/main.py")

    print(f"Подключаюсь к БД {DB_CONFIG['user']}@{DB_CONFIG['host']}:"
          f"{DB_CONFIG['port']}/{DB_CONFIG['dbname']}...")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
    except psycopg2.OperationalError as exc:
        sys.exit(f"Не удалось подключиться к БД (запускали etl/main.py?): {exc}")
    try:
        top5 = queries.top5_overrun_detailed(conn)
        relative = queries.top5_overrun_relative(conn)
        cost = queries.overrun_cost(conn)
        stages = queries.stage_estimate_ratio(conn)
        balance = queries.buffer_balance(conn)
        leak = queries.buffer_leak(conn)
        cats = queries.category_estimate_ratio(conn)
        concentration = queries.budget_concentration(conn)
        pm_rows = queries.top3_pm_budget(conn)
    finally:
        conn.close()

    pm_top3 = [(r[0], round(float(r[1]) / 1e9, 2)) for r in pm_rows[:3]]

    print("Собираю слайды...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1(prs, KPI, cost)
    slide_2(prs, top5, relative, cost)
    slide_3(prs, stages)
    slide_4(prs, balance, stages, leak)
    slide_5(prs, cats)
    slide_6(prs, concentration, pm_top3)
    slide_7(prs)

    prs.save(OUT_FILE)
    print(f"Готово: {len(prs.slides._sldIdLst)} слайдов -> {OUT_FILE}")


if __name__ == "__main__":
    main()
